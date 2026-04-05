"""
Generate training/slides.pptx — DreamAirlines Data Engineering Training.
Run:  python training/generate_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── palette ──────────────────────────────────────────────────────────────────
PURPLE       = RGBColor(0x6D, 0x28, 0xD9)
PURPLE_LIGHT = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_DARK  = RGBColor(0x3B, 0x07, 0x64)
PURPLE_MUTED = RGBColor(0xED, 0xE9, 0xFE)
BLACK        = RGBColor(0x0F, 0x0F, 0x0F)
CHARCOAL     = RGBColor(0x1F, 0x1F, 0x2E)
GRAY         = RGBColor(0x6B, 0x72, 0x80)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG     = RGBColor(0xF3, 0xF0, 0xFF)

# slide dimensions: 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

# ── helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    fill_solid(shape, color)
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height,
                text, font_size, bold=False, italic=False,
                color=BLACK, align=PP_ALIGN.LEFT,
                font_name="Calibri", wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font_name
    if line_spacing:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txBox


def add_multiline_textbox(slide, left, top, width, height,
                          lines,           # list of (text, font_size, bold, color)
                          line_spacing_pt=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, font_size, bold, color) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if line_spacing_pt:
            from pptx.oxml.ns import qn
            import lxml.etree as etree
            pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing_pt * 100)))
        run = p.add_run()
        run.text = text
        run.font.size      = Pt(font_size)
        run.font.bold      = bold
        run.font.color.rgb = color
        run.font.name      = "Calibri"
    return txBox


def add_code_box(slide, left, top, width, height, code_text, font_size=9):
    """Monospaced code block on a light purple background."""
    bg = add_rect(slide, left, top, width, height, LIGHT_BG)
    # left accent bar
    add_rect(slide, left, top, Inches(0.06), height, PURPLE)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.12), top + Inches(0.1),
        width - Inches(0.18), height - Inches(0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in code_text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size      = Pt(font_size)
        run.font.color.rgb = CHARCOAL
        run.font.name      = "Courier New"
    return txBox


def slide_header(slide, tag, title):
    """Dark-purple header band across the top of content slides."""
    add_rect(slide, 0, 0, W, Inches(0.75), PURPLE_DARK)
    add_textbox(slide, Inches(0.5), Inches(0.13), Inches(2), Inches(0.35),
                tag.upper(), 8, bold=True, color=PURPLE_LIGHT, font_name="Calibri")
    add_textbox(slide, Inches(2.3), Inches(0.13), Inches(10), Inches(0.35),
                title, 16, bold=True, color=WHITE, font_name="Calibri")


def section_label(slide, left, top, text):
    add_textbox(slide, left, top, Inches(5.5), Inches(0.3),
                text.upper(), 8, bold=True, color=PURPLE, font_name="Calibri")
    # underline via thin rect
    add_rect(slide, left, top + Inches(0.28), Inches(1.5), Emu(18000), PURPLE_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════

def slide1_title(prs):
    slide = blank_slide(prs)

    # white background
    add_rect(slide, 0, 0, W, H, WHITE)

    # left purple accent strip
    add_rect(slide, 0, 0, Inches(0.35), H, PURPLE_DARK)

    # kicker
    add_textbox(slide, Inches(0.65), Inches(1.6), Inches(8), Inches(0.4),
                "DATA ENGINEERING TRAINING", 11, bold=True, color=PURPLE,
                font_name="Calibri")

    # main title — two lines
    add_textbox(slide, Inches(0.65), Inches(2.05), Inches(8), Inches(1.1),
                "Dream", 60, bold=True, color=BLACK, font_name="Calibri")
    add_textbox(slide, Inches(2.85), Inches(2.05), Inches(8), Inches(1.1),
                "Airlines", 60, bold=True, color=PURPLE, font_name="Calibri")

    # subtitle
    add_textbox(slide, Inches(0.65), Inches(3.1), Inches(9), Inches(0.4),
                "Your first real pipeline.  Four notebooks.  One mission.", 14,
                italic=True, color=GRAY, font_name="Calibri")

    # mission box background
    mission_top = Inches(3.75)
    mission_h   = Inches(2.7)
    add_rect(slide, Inches(0.65), mission_top, Inches(9.8), mission_h, CHARCOAL)
    add_rect(slide, Inches(0.65), mission_top, Inches(0.06), mission_h, PURPLE_LIGHT)

    # mission label
    add_textbox(slide, Inches(0.85), mission_top + Inches(0.15), Inches(9), Inches(0.3),
                "// INCOMING TRANSMISSION", 9, bold=True, color=PURPLE_LIGHT,
                font_name="Courier New")

    mission_text = (
        "Good morning, Agent.\n\n"
        "Your target: DreamAirlines — a Portuguese travel agency with three offices\n"
        "and data scattered across 14 source tables.  Nobody has ever joined them together.\n"
        "The analytics team is flying blind.\n\n"
        "Your mission, should you choose to accept it, is to ingest the raw data,\n"
        "expose its flaws, and deliver a clean revenue pipeline the business can actually trust.\n\n"
        "The Spark cluster is standing by."
    )
    add_textbox(slide, Inches(0.85), mission_top + Inches(0.5),
                Inches(9.4), Inches(2.0),
                mission_text, 10, color=WHITE, font_name="Calibri")

    # footer
    add_rect(slide, 0, H - Inches(0.45), W, Inches(0.45), PURPLE_DARK)
    add_textbox(slide, Inches(0.65), H - Inches(0.38), Inches(10), Inches(0.3),
                "Databricks   ·   PySpark   ·   Delta Lake   ·   MLflow",
                10, color=PURPLE_LIGHT, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROJECT STRUCTURE
# ═══════════════════════════════════════════════════════════════════════════════

def slide2_structure(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, W, H, WHITE)
    slide_header(slide, "The Boilerplate", "How every ETL job in this repo is structured")

    top = Inches(0.95)

    # ── LEFT COLUMN ──────────────────────────────────────────────
    lx = Inches(0.4)
    col_w = Inches(5.8)

    section_label(slide, lx, top, "Repository layout")
    repo_tree = (
        "pyspark_etl_boilerplate/\n"
        "├── main.py              # entry point\n"
        "├── config/\n"
        "│   ├── dev/             # per-environment\n"
        "│   ├── test/\n"
        "│   ├── qa/\n"
        "│   └── prod/\n"
        "├── etl/\n"
        "│   ├── etl_interface.py # abstract base\n"
        "│   ├── etl_config.py    # config loader\n"
        "│   └── my_etl_job.py    # sample job\n"
        "└── utils/\n"
        "    ├── config.py        # SparkSession\n"
        "    ├── toolbox.py\n"
        "    └── ..."
    )
    add_code_box(slide, lx, top + Inches(0.32), col_w, Inches(2.55), repo_tree, font_size=8.5)

    section_label(slide, lx, top + Inches(3.05), "Config sections")
    pills = ["[SPARK_CONTEXT]", "[COS]", "[DATABASE]", "[SPARK]", "[KAFKA]"]
    px = lx
    for pill in pills:
        pw = Inches(len(pill) * 0.085 + 0.2)
        bg = add_rect(slide, px, top + Inches(3.38), pw, Inches(0.28), PURPLE_MUTED)
        add_textbox(slide, px + Inches(0.06), top + Inches(3.4), pw, Inches(0.25),
                    pill, 8, bold=True, color=PURPLE_DARK, font_name="Courier New")
        px += pw + Inches(0.08)

    note = ("All Spark session properties are sourced from [SPARK_CONTEXT].\n"
            "Placeholders like {bucket} are resolved from [COS] at read time.")
    add_textbox(slide, lx, top + Inches(3.78), col_w, Inches(0.5),
                note, 9, color=GRAY, font_name="Calibri")

    # ── RIGHT COLUMN ─────────────────────────────────────────────
    rx = Inches(6.6)

    section_label(slide, rx, top, "Execution flow")
    add_textbox(slide, rx, top + Inches(0.33), Inches(6.4), Inches(0.25),
                "spark-submit main.py  --className …  --confPath …",
                8.5, color=GRAY, align=PP_ALIGN.CENTER, font_name="Courier New")

    # flow boxes
    flow_items = ["__init__", "extract()", "transform()", "load()"]
    flow_top = top + Inches(0.65)
    box_w = Inches(1.12)
    box_h = Inches(0.38)
    gap   = Inches(0.18)
    total_w = len(flow_items) * box_w + (len(flow_items) - 1) * (gap + Inches(0.3))
    fx = rx + (Inches(6.4) - total_w) / 2

    for i, label in enumerate(flow_items):
        bg_color = PURPLE if label == "transform()" else PURPLE_DARK
        add_rect(slide, fx, flow_top, box_w, box_h, bg_color)
        add_textbox(slide, fx, flow_top + Inches(0.06), box_w, Inches(0.28),
                    label, 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                    font_name="Calibri")
        fx += box_w
        if i < len(flow_items) - 1:
            add_textbox(slide, fx + Inches(0.03), flow_top + Inches(0.06),
                        Inches(0.3), Inches(0.28),
                        "→", 12, color=PURPLE_LIGHT, align=PP_ALIGN.CENTER,
                        font_name="Calibri")
            fx += gap + Inches(0.3)

    add_textbox(slide, rx, flow_top + Inches(0.44), Inches(6.4), Inches(0.25),
                "↑  YOUR WORK LIVES HERE", 8, bold=True, color=PURPLE,
                align=PP_ALIGN.CENTER, font_name="Calibri")

    section_label(slide, rx, top + Inches(1.3), "Every job follows the same contract")
    contract_code = (
        "class MyETLJob(ETLInterface):\n"
        "\n"
        "    def extract(self):\n"
        "        # read from source\n"
        "\n"
        "    def transform(self):\n"
        "        # clean, join, aggregate\n"
        "\n"
        "    def load(self):\n"
        "        # write output\n"
        "\n"
        "    def run(self):\n"
        "        with mlflow.start_run():\n"
        "            self.extract()\n"
        "            self.transform()\n"
        "            self.load()"
    )
    add_code_box(slide, rx, top + Inches(1.62), Inches(6.4), Inches(2.9), contract_code)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE ASSIGNMENT
# ═══════════════════════════════════════════════════════════════════════════════

def slide3_assignment(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, W, H, WHITE)
    slide_header(slide, "The Assignment", "Four notebooks — one pipeline")

    top = Inches(0.95)
    lx  = Inches(0.4)
    rx  = Inches(6.8)
    col_w = Inches(6.0)

    # ── LEFT: notebook steps ──────────────────────────────────────
    section_label(slide, lx, top, "Run in order")

    steps = [
        ("00", "Setup — load the data",
         "14 CSV files → raw.* Delta tables.  One command, done."),
        ("01", "Explore — profile the data",
         "Find nulls, broken FK, integer dates, audit noise.  Document every finding."),
        ("02", "Build — complete the ETL job",
         "Fill in 4 TODOs inside transform().  Everything else is wired up."),
        ("04", "Upgrade — rewrite SQL as DataFrame API",
         "Convert a working SQL job to the DataFrame API.  Then watch MLflow catch failures."),
    ]

    sy = top + Inches(0.35)
    for (num, title, desc) in steps:
        # circle badge
        add_rect(slide, lx, sy, Inches(0.42), Inches(0.42), PURPLE_DARK)
        add_textbox(slide, lx, sy + Inches(0.04), Inches(0.42), Inches(0.32),
                    num, 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                    font_name="Calibri")
        add_textbox(slide, lx + Inches(0.52), sy, Inches(5.6), Inches(0.25),
                    title, 11, bold=True, color=BLACK, font_name="Calibri")
        add_textbox(slide, lx + Inches(0.52), sy + Inches(0.25), Inches(5.6), Inches(0.35),
                    desc, 9, color=GRAY, font_name="Calibri")
        sy += Inches(0.75)

    # ── RIGHT: target schema + quality issues ─────────────────────
    section_label(slide, rx, top, "What you will build")
    schema = (
        "analytics.booking_revenue_summary\n"
        "──────────────────────────────────\n"
        "booking_date       DATE\n"
        "travel_agency      STRING\n"
        "destination        STRING\n"
        "travel_class       STRING\n"
        "customer_type      STRING\n"
        "num_bookings       LONG\n"
        "travel_revenue     DOUBLE\n"
        "hotel_revenue      DOUBLE\n"
        "car_revenue        DOUBLE\n"
        "insurance_revenue  DOUBLE\n"
        "total_revenue      DOUBLE\n"
        "revenue_tier       STRING   # High / Medium / Low"
    )
    add_code_box(slide, rx, top + Inches(0.32), Inches(6.4), Inches(2.35), schema, font_size=8.5)

    section_label(slide, rx, top + Inches(2.85), "Data quality issues to find first")
    issues = [
        ("Q1", "Nullable FK in fact_car_renting",    "LEFT JOIN"),
        ("Q2", "Null amounts across all facts",       "COALESCE → 0"),
        ("Q3", "Dates as integers (YYYYMMDD)",        "join dim_date"),
        ("Q4", "Customer join is on name, not PK",    "DSC_Customer"),
    ]
    iy = top + Inches(3.18)
    # header row
    for (cx, label) in [(rx, "#"), (rx + Inches(0.5), "Issue"), (rx + Inches(4.2), "Fix")]:
        add_textbox(slide, cx, iy, Inches(3.5), Inches(0.25),
                    label, 9, bold=True, color=PURPLE_DARK, font_name="Calibri")
    add_rect(slide, rx, iy + Inches(0.26), Inches(6.4), Emu(12000), PURPLE_MUTED)
    iy += Inches(0.32)
    for (num, issue, fix) in issues:
        add_textbox(slide, rx, iy, Inches(0.45), Inches(0.28),
                    num, 9, bold=True, color=PURPLE, font_name="Calibri")
        add_textbox(slide, rx + Inches(0.5), iy, Inches(3.6), Inches(0.28),
                    issue, 9, color=BLACK, font_name="Calibri")
        # fix badge
        add_rect(slide, rx + Inches(4.15), iy, Inches(1.8), Inches(0.26), PURPLE_MUTED)
        add_textbox(slide, rx + Inches(4.2), iy + Inches(0.02), Inches(1.75), Inches(0.24),
                    fix, 8, bold=True, color=PURPLE_DARK, font_name="Calibri")
        iy += Inches(0.33)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTIONS & TAKEAWAYS
# ═══════════════════════════════════════════════════════════════════════════════

def slide4_solutions(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, W, H, WHITE)
    slide_header(slide, "Solutions & Takeaways", "What the completed pipeline looks like")

    top = Inches(0.95)
    lx  = Inches(0.4)
    rx  = Inches(6.8)

    # ── LEFT: 4 TODOs solved ──────────────────────────────────────
    section_label(slide, lx, top, "The four TODOs — solved")

    todos = [
        ("T1", "Join chain",
         ".join(dim_date, …, \"inner\")  then filter,\n"
         "then .join(dim_customer, …),\n"
         "then four \"left\" joins for the facts"),
        ("T2", "COALESCE amounts",
         "F.coalesce(F.col(\"AMT_Travel\"), F.lit(0))\n"
         "— applied after joins, before aggregation"),
        ("T3", "Aggregate",
         ".groupBy(date, agency, country, class, cust_type)\n"
         ".agg(F.countDistinct(…), F.round(F.sum(…), 2), …)"),
        ("T4", "Revenue tier",
         "F.when(col >= 2000, \"High\")\n"
         " .when(col >= 500,  \"Medium\")\n"
         " .otherwise(\"Low\")"),
    ]

    sy = top + Inches(0.35)
    for (num, title, desc) in todos:
        add_rect(slide, lx, sy, Inches(0.42), Inches(0.42), PURPLE_DARK)
        add_textbox(slide, lx, sy + Inches(0.04), Inches(0.42), Inches(0.32),
                    num, 8, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                    font_name="Calibri")
        add_textbox(slide, lx + Inches(0.52), sy, Inches(5.6), Inches(0.25),
                    title, 11, bold=True, color=BLACK, font_name="Calibri")
        add_code_box(slide, lx + Inches(0.52), sy + Inches(0.27),
                     Inches(5.6), Inches(0.55), desc, font_size=8)
        sy += Inches(0.98)

    # ── RIGHT: API vs SQL table + MLflow ─────────────────────────
    section_label(slide, rx, top, "DataFrame API vs SQL — the key difference")

    comparisons = [
        ("",               "SQL",                "DataFrame API"),
        ("Multi-step logic","CTEs in one string", "Named variables — inspect any step"),
        ("Reuse a column",  "Must repeat formula","Compute once, reference by name"),
        ("Debug mid-pipe",  "Rewrite whole query","display(step) anywhere"),
        ("IDE support",     "Plain string",       "Autocomplete + type hints"),
    ]
    ty = top + Inches(0.33)
    col_widths = [Inches(1.55), Inches(2.1), Inches(2.55)]
    for i, (c1, c2, c3) in enumerate(comparisons):
        row_color = PURPLE_DARK if i == 0 else (LIGHT_BG if i % 2 == 0 else WHITE)
        add_rect(slide, rx, ty, sum(col_widths), Inches(0.3), row_color)
        txt_color = WHITE if i == 0 else BLACK
        for j, (cx_off, text, cw) in enumerate(zip(
            [0, col_widths[0], col_widths[0] + col_widths[1]],
            [c1, c2, c3], col_widths
        )):
            add_textbox(slide, rx + cx_off + Inches(0.06), ty + Inches(0.04),
                        cw - Inches(0.08), Inches(0.25),
                        text, 8.5, bold=(i == 0), color=txt_color, font_name="Calibri")
        ty += Inches(0.3)

    section_label(slide, rx, top + Inches(2.4), "MLflow — why every run gets logged")
    mlflow_code = (
        "# Find gaps in the daily pipeline\n"
        "runs = mlflow.search_runs(\n"
        "    filter_string=\"tags.job = 'DailyBookingRevenueJob'\"\n"
        ")\n"
        "failed = runs[runs[\"status\"] != \"FINISHED\"]\n"
        "# → reruns only the missing dates"
    )
    add_code_box(slide, rx, top + Inches(2.72), Inches(6.4), Inches(1.35),
                 mlflow_code, font_size=8.5)

    note = ("Log every run — even failures.  Query the gaps.  Backfill with one loop.\n"
            "This is the operational backbone of any reliable daily pipeline.")
    add_textbox(slide, rx, top + Inches(4.15), Inches(6.4), Inches(0.5),
                note, 9, color=GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import os
    prs = new_prs()

    slide1_title(prs)
    slide2_structure(prs)
    slide3_assignment(prs)
    slide4_solutions(prs)

    out = os.path.join(os.path.dirname(__file__), "slides.pptx")
    prs.save(out)
    print(f"Saved: {out}")
